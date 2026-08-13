"""
MLVerse X — RAG System
Handles document ingestion, chunking, embedding, and retrieval.
"""
import io
import uuid
import logging
from typing import List, Optional, Dict, Any, Tuple
from pathlib import Path

from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.config import Settings as ChromaSettings

from app.core.config import settings

logger = logging.getLogger(__name__)

# ─── Embedding Model (local, no cloud) ─────────────────────────────────────────
_embedding_model: Optional[SentenceTransformer] = None

def get_embedding_model() -> SentenceTransformer:
    global _embedding_model
    if _embedding_model is None:
        logger.info(f"Loading embedding model: {settings.EMBEDDING_MODEL}")
        _embedding_model = SentenceTransformer(
            settings.EMBEDDING_MODEL,
            device=settings.EMBEDDING_DEVICE,
        )
    return _embedding_model


# ─── ChromaDB Client ───────────────────────────────────────────────────────────
class AsyncLocalCollection:
    def __init__(self, collection):
        self._collection = collection

    async def add(self, ids, documents, embeddings, metadatas):
        self._collection.add(ids=ids, documents=documents, embeddings=embeddings, metadatas=metadatas)

    async def query(self, query_embeddings, n_results, include):
        return self._collection.query(query_embeddings=query_embeddings, n_results=n_results, include=include)

    async def count(self):
        return self._collection.count()

class AsyncLocalChromaClient:
    def __init__(self, client):
        self._client = client

    async def get_or_create_collection(self, name, metadata=None):
        col = self._client.get_or_create_collection(name=name, metadata=metadata)
        return AsyncLocalCollection(col)

    async def get_collection(self, name):
        col = self._client.get_collection(name=name)
        return AsyncLocalCollection(col)

    async def delete_collection(self, name):
        self._client.delete_collection(name=name)

_chroma_client = None

async def get_chroma_client():
    global _chroma_client
    if _chroma_client is None:
        remote_ok = False
        try:
            import socket
            sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
            sock.settimeout(0.5)
            result = sock.connect_ex((settings.CHROMA_HOST, settings.CHROMA_PORT))
            sock.close()
            if result == 0:
                remote_ok = True
        except Exception:
            pass

        if remote_ok:
            try:
                _chroma_client = await chromadb.AsyncHttpClient(
                    host=settings.CHROMA_HOST,
                    port=settings.CHROMA_PORT,
                )
                logger.info(f"Successfully connected to ChromaDB server at {settings.CHROMA_HOST}:{settings.CHROMA_PORT}")
            except Exception as e:
                remote_ok = False

        if not remote_ok:
            db_path = Path(__file__).resolve().parents[3] / "chroma_db"
            db_path.mkdir(parents=True, exist_ok=True)
            logger.warning(f"ChromaDB server at {settings.CHROMA_HOST}:{settings.CHROMA_PORT} is unreachable. Falling back to local persistent client at {db_path}.")
            sync_client = chromadb.PersistentClient(path=str(db_path))
            _chroma_client = AsyncLocalChromaClient(sync_client)

    return _chroma_client


# ─── Document Parsers ──────────────────────────────────────────────────────────
class DocumentParser:
    """Parse various file formats into plain text."""

    @staticmethod
    def parse_pdf(file_bytes: bytes) -> str:
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes))
            page_texts = []
            for idx, page in enumerate(reader.pages):
                txt = page.extract_text() or ""
                if txt.strip():
                    page_texts.append(f"--- [Page {idx + 1}] ---\n{txt.strip()}")
            return "\n\n".join(page_texts).strip()
        except Exception as e:
            logger.warning(f"Error parsing PDF with pypdf: {e}")
            try:
                return file_bytes.decode("utf-8", errors="ignore").strip()
            except Exception:
                return ""

    @staticmethod
    def parse_docx(file_bytes: bytes) -> str:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    @staticmethod
    def parse_xlsx(file_bytes: bytes) -> str:
        import pandas as pd
        df = pd.read_excel(io.BytesIO(file_bytes))
        return df.to_string(index=False)

    @staticmethod
    def parse_csv(file_bytes: bytes) -> str:
        import pandas as pd
        df = pd.read_csv(io.BytesIO(file_bytes))
        return df.to_string(index=False)

    @staticmethod
    def parse_pptx(file_bytes: bytes) -> str:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n\n".join(texts)

    @staticmethod
    def parse_image_ocr(file_bytes: bytes) -> str:
        import pytesseract
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        return pytesseract.image_to_string(img)

    @staticmethod
    def parse_audio(file_bytes: bytes, filename: str) -> str:
        import whisper
        import tempfile, os
        suffix = Path(filename).suffix
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as f:
            f.write(file_bytes)
            tmp_path = f.name
        try:
            model = whisper.load_model("base")
            result = model.transcribe(tmp_path)
            return result["text"]
        finally:
            os.unlink(tmp_path)

    @classmethod
    def parse(cls, file_bytes: bytes, filename: str, file_type: str) -> str:
        ext = file_type.lower().strip(".")
        if ext == "pdf":
            return cls.parse_pdf(file_bytes)
        elif ext in ("docx", "doc"):
            return cls.parse_docx(file_bytes)
        elif ext in ("xlsx", "xls"):
            return cls.parse_xlsx(file_bytes)
        elif ext == "csv":
            return cls.parse_csv(file_bytes)
        elif ext == "pptx":
            return cls.parse_pptx(file_bytes)
        elif ext in ("txt", "md", "py", "js", "ts", "java", "cpp", "c", "html", "sql"):
            return file_bytes.decode("utf-8", errors="ignore")
        elif ext in ("png", "jpg", "jpeg", "bmp", "tiff"):
            return cls.parse_image_ocr(file_bytes)
        elif ext in ("mp3", "wav", "m4a", "ogg", "flac"):
            return cls.parse_audio(file_bytes, filename)
        else:
            try:
                return file_bytes.decode("utf-8", errors="ignore")
            except Exception:
                return ""


# ─── Chunker ───────────────────────────────────────────────────────────────────
class SemanticChunker:
    """Split text into overlapping semantic chunks."""

    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def chunk(self, text: str, document_name: str) -> List[Dict[str, Any]]:
        words = text.split()
        chunks = []
        i = 0
        chunk_index = 0

        while i < len(words):
            chunk_words = words[i : i + self.chunk_size]
            chunk_text = " ".join(chunk_words)

            if len(chunk_text.strip()) > 20:  # skip tiny chunks
                chunks.append({
                    "id": f"{document_name}_{chunk_index}",
                    "text": chunk_text,
                    "metadata": {
                        "document": document_name,
                        "chunk_index": chunk_index,
                        "word_count": len(chunk_words),
                    },
                })
                chunk_index += 1

            i += self.chunk_size - self.chunk_overlap

        return chunks


# ─── RAG Pipeline ──────────────────────────────────────────────────────────────
class RAGPipeline:
    """Full RAG pipeline: ingest → embed → store → retrieve → generate."""

    def __init__(self):
        self.chunker = SemanticChunker(chunk_size=500, chunk_overlap=50)

    async def ingest_document(
        self,
        file_bytes: bytes,
        filename: str,
        file_type: str,
        collection_id: str,
        user_id: int,
    ) -> Dict[str, Any]:
        """Parse, chunk, embed, and store a document in ChromaDB."""
        # Parse document
        text = DocumentParser.parse(file_bytes, filename, file_type)
        if not text.strip():
            raise ValueError(f"Could not extract text from {filename}")

        # Chunk
        chunks = self.chunker.chunk(text, filename)
        if not chunks:
            raise ValueError("No meaningful chunks extracted from document")

        # Embed
        model = get_embedding_model()
        texts = [c["text"] for c in chunks]
        embeddings = model.encode(texts, batch_size=32, show_progress_bar=False).tolist()

        # Store in ChromaDB
        client = await get_chroma_client()
        collection = await client.get_or_create_collection(
            name=collection_id,
            metadata={"user_id": str(user_id)},
        )

        await collection.add(
            ids=[c["id"] for c in chunks],
            documents=texts,
            embeddings=embeddings,
            metadatas=[c["metadata"] for c in chunks],
        )

        return {
            "filename": filename,
            "num_chunks": len(chunks),
            "collection_id": collection_id,
            "char_count": len(text),
        }

    async def retrieve(
        self,
        query: str,
        collection_id: str,
        top_k: int = 5,
    ) -> List[Dict[str, Any]]:
        """Retrieve relevant chunks for a query."""
        model = get_embedding_model()
        query_embedding = model.encode([query], show_progress_bar=False).tolist()[0]

        client = await get_chroma_client()
        try:
            collection = await client.get_collection(name=collection_id)
        except Exception:
            return []

        results = await collection.query(
            query_embeddings=[query_embedding],
            n_results=min(top_k, await collection.count()),
            include=["documents", "metadatas", "distances"],
        )

        chunks = []
        for i, (doc, meta, dist) in enumerate(zip(
            results["documents"][0],
            results["metadatas"][0],
            results["distances"][0],
        )):
            chunks.append({
                "text": doc,
                "metadata": meta,
                "relevance_score": float(1 - dist),  # convert distance to similarity
                "rank": i + 1,
            })

        return chunks

    async def answer(
        self,
        query: str,
        collection_ids: List[str],
        model: str = None,
        top_k: int = 5,
    ) -> Dict[str, Any]:
        """Full RAG: retrieve context + generate answer via Ollama."""
        from ai.ollama.client import get_ollama_client, RAG_SYSTEM_PROMPT

        # Retrieve from all collections
        all_chunks = []
        for cid in collection_ids:
            chunks = await self.retrieve(query, cid, top_k=top_k)
            all_chunks.extend(chunks)

        # Sort by relevance
        all_chunks.sort(key=lambda x: x["relevance_score"], reverse=True)
        top_chunks = all_chunks[:top_k]

        # Build context
        context = "\n\n---\n\n".join([
            f"[Source: {c['metadata'].get('document', 'unknown')} | Chunk {c['metadata'].get('chunk_index', 0)}]\n{c['text']}"
            for c in top_chunks
        ])

        # Generate answer
        client = get_ollama_client()
        system = RAG_SYSTEM_PROMPT.format(context=context)
        messages = [{"role": "user", "content": query}]

        answer_parts = []
        async for chunk in client.chat(
            messages=messages,
            model=model,
            stream=True,
            system_prompt=system,
        ):
            answer_parts.append(chunk)

        answer = "".join(answer_parts)

        # Citations
        citations = [
            {
                "document": c["metadata"].get("document", ""),
                "chunk_index": c["metadata"].get("chunk_index", 0),
                "relevance": round(c["relevance_score"], 3),
                "text_preview": c["text"][:200] + "..." if len(c["text"]) > 200 else c["text"],
            }
            for c in top_chunks
        ]

        return {
            "query": query,
            "answer": answer,
            "citations": citations,
            "context_chunks_used": len(top_chunks),
        }

    async def explain_pdf(
        self,
        file_bytes: bytes,
        filename: str,
        model: str = None,
    ) -> Dict[str, Any]:
        """Extract text from PDF and generate an executive summary, key findings, and topic breakdown."""
        text = DocumentParser.parse_pdf(file_bytes)
        if not text.strip():
            raise ValueError(f"Could not extract readable text from PDF '{filename}'")

        char_count = len(text)
        word_count = len(text.split())

        # Generate overview prompt
        prompt = (
            f"Please read the following PDF document content from '{filename}' and provide a comprehensive, structured explanation:\n"
            f"1. Executive Summary (3-4 sentences)\n"
            f"2. Key Takeaways & Main Topics Covered\n"
            f"3. Important Metrics, Data Points, or Conclusions (if present)\n"
            f"4. 3 Suggested Questions to ask about this PDF\n\n"
            f"PDF Content Excerpt (First 4000 characters):\n{text[:4000]}"
        )

        from ai.ollama.client import get_ollama_client
        client = get_ollama_client()
        explanation = await client.generate(prompt=prompt, model=model)

        return {
            "filename": filename,
            "word_count": word_count,
            "char_count": char_count,
            "explanation": explanation,
            "preview_snippet": text[:500] + ("..." if len(text) > 500 else ""),
        }

    async def delete_collection(self, collection_id: str):
        client = await get_chroma_client()
        try:
            await client.delete_collection(collection_id)
        except Exception as e:
            logger.warning(f"Failed to delete collection {collection_id}: {e}")

